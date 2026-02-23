#!/usr/bin/env python3
"""Generate A3 landscape single-page PPTX summary table."""

from pptx import Presentation
from pptx.util import Mm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import os

# ── Constants ──────────────────────────────────────────────────────────────
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEADER_BG = NAVY
LIGHT_BG = RGBColor(0xF7, 0xF9, 0xFC)

OUTPUT_PATH = os.path.join(
    os.path.dirname(__file__),
    "public", "downloads", "summary-a3.pptx",
)

# A3 landscape dimensions
SLIDE_W = Mm(420)
SLIDE_H = Mm(297)

# ── Table data ─────────────────────────────────────────────────────────────
HEADERS = [
    "단계",
    "산업 시대",
    "AI 시대:\n왜 다른가",
    "\u2776 주체 정의",
    "\u2777 관계 정의",
    "\u2778 질서 설계",
    "목표",
]

# Each cell: list of (text, bold) tuples
# Using bullet points with bold keywords

ROWS = [
    # Row 1: 문제 인식
    [
        # 단계
        [("문제 인식", True), ("\n기술 혁명은 왜\n위험한가", False)],
        # 산업 시대
        [
            ("\u2022 ", False), ("5차례 기술 혁명", True), (" 모두\n  제도 재편 없이 황금기\n  도달한 사례 없음", False),
            ("\n\u2022 ", False), ("선제적 설계", True), ("가\n  격변보다 사회적 비용을\n  현저히 줄임", False),
            ("\n\u2022 ", False), ("진단은 풍부", True), ("하나\n  황금기로 가는 ", False), ("구체적\n  경로", True), ("가 부재", False),
        ],
        # AI 시대
        [
            ("\u2022 ", False), ("사후 교정", True), ("의 비용과\n  위험이 ", False), ("질적으로 변화", True),
            ("\n\u2022 질서의 일부를 ", False), ("미리\n  설계", True), ("할 필요", False),
            ("\n\u2022 질서 자체가 아닌 ", False), ("질서가\n  작동할 조건", True), ("의 설계", False),
        ],
        # 주체 정의
        [("(이 단계에서\n도출되는 과제)", False)],
        # 관계 정의
        [("(이 단계에서\n도출되는 과제)", False)],
        # 질서 설계
        [("(이 단계에서\n도출되는 과제)", False)],
        # 목표
        [("(이 단계에서\n도출되는 과제)", False)],
    ],
    # Row 2: 역사적 분석
    [
        # 단계
        [("역사적 분석", True), ("\n산업 시대의\n질서 형성", False)],
        # 산업 시대
        [
            ("\u2022 ", False), ("시장이 먼저, 법이 나중", True),
            ("\n  \u2013 보이지 않는 손이 질서\n    형성", False),
            ("\n\u2022 ", False), ("사후 제도화 가능", True), ("한 이유:", False),
            ("\n  \u2013 기계는 인간이 끄면\n    멈춤", False),
            ("\n  \u2013 폴라니의 ", False), ("이중 운동", True), (":\n    반격할 시간이 존재", False),
            ("\n\u2022 ", False), ("질서 \u2192 거버넌스", True), (" 순서", False),
            ("\n  \u2013 법인\u00b7노동법\u00b7복지국가\n    순차 형성", False),
        ],
        # AI 시대
        [("-", False)],
        # 주체 정의
        [("-", False)],
        # 관계 정의
        [("-", False)],
        # 질서 설계
        [("-", False)],
        # 목표
        [("-", False)],
    ],
    # Row 3: 현재 진단
    [
        # 단계
        [("현재 진단", True), ("\nAI 시대는\n왜 다른가", False)],
        # 산업 시대
        [("-", False)],
        # AI 시대
        [
            ("\u2022 ", False), ("책임 귀속 공백", True),
            ("\n  \u2013 AI 행위의 법적 귀속처\n    부재", False),
            ("\n  \u2013 불확실성 비용 누적", False),
            ("\n\u2022 ", False), ("사후 통제 한계", True),
            ("\n  \u2013 통제를 벗어날 가능성이\n    있는 최초의 범용 기술", False),
            ("\n  \u2013 되돌릴 수 없는 결과 위험", False),
            ("\n\u2022 ", False), ("속도 격차", True),
            ("\n  \u2013 사회 반격 시간 부족", False),
            ("\n  \u2013 확산 후 통제 어려움", False),
        ],
        # 주체 정의
        [("-", False)],
        # 관계 정의
        [("-", False)],
        # 질서 설계
        [("-", False)],
        # 목표
        [("-", False)],
    ],
    # Row 4: 해법 제시 — 3단계 경로
    [
        # 단계
        [("해법 제시", True), ("\n세 단계의 경로", False)],
        # 산업 시대
        [
            ("\u2022 법인 범주 창설(1844)", False),
            ("\n  \u2192 이후 모든 제도의 토대", False),
            ("\n\u2022 가격\u00b7계약 \u2192 관계 형성", False),
            ("\n  \u2192 재산권\u00b7노동법 정비", False),
            ("\n\u2022 3주체 역할 확립", False),
            ("\n  (법인\u00b7자연인\u00b7국가)", False),
        ],
        # AI 시대
        [
            ("\u2022 질서 자체가 아닌\n  ", False), ("조건의 설계", True),
            ("\n\u2022 브레튼우즈\u00b7EU 단일시장\n  처럼 ", False), ("틀", True), ("을 만드는 것", False),
            ("\n\u2022 하이에크 반론 수용:", False),
            ("\n  보이지 않는 손의 ", False), ("울타리", True),
        ],
        # 주체 정의
        [
            ("\u2022 ", False), ("전자인", True), (" 법적 범주 창설", False),
            ("\n\u2022 자율성 ", False), ("스펙트럼", True), (" 설계:", False),
            ("\n  ", False), ("도구 \u2192 대리인 \u2192\n  준자율 \u2192 자율", True),
            ("\n\u2022 각 단계별:", False),
            ("\n  \u2013 자율성 범위 차등", False),
            ("\n  \u2013 책임 귀속 구조 차등", False),
            ("\n  \u2013 인간 개입 수준 차등", False),
            ("\n\u2022 러셀의 ", False), ("번복 가능성", True),
            (" 원칙이\n  전체 관통", False),
        ],
        # 관계 정의
        [
            ("\u2022 3당사자 관계 규율\n  (자연인\u00b7법인\u00b7전자인)", False),
            ("\n\u2022 4대 영역:", False),
            ("\n  \u2013 ", False), ("데이터 소유권", True),
            ("\n  \u2013 ", False), ("인간-AI 책임 배분", True),
            ("\n  \u2013 ", False), ("알고리즘 권력", True), (" 규율", False),
            ("\n  \u2013 ", False), ("디지털 시민권", True),
            ("\n\u2022 산업별 파편 \u2192\n  ", False), ("일반 책임 구조", True), (" 초안", False),
        ],
        # 질서 설계
        [
            ("\u2022 ", False), ("4 경제주체", True), (" 역할 재정의", False),
            ("\n  \u2013 전자인: 경제행위 범위", False),
            ("\n  \u2013 가계: 과실 수령 경로", False),
            ("\n  \u2013 기업: AI 책임 범위", False),
            ("\n  \u2013 정부: 초국경 거버넌스", False),
            ("\n\u2022 ", False), ("분배 4후보", True), (" 통합:", False),
            ("\n  기본소득\u00b7데이터배당\u00b7", False),
            ("\n  공공AI인프라\u00b7역량보장", False),
            ("\n\u2022 ", False), ("시민참여형 거버넌스", True),
        ],
        # 목표
        [
            ("\u2022 ", False), ("노동 중심 \u2192 생명 중심", True),
            ("\n  가치 전환", False),
            ("\n\u2022 인간의 가치 =\n  ", False), ("생산성이 아닌\n  존엄 그 자체", True),
            ("\n\u2022 역량 접근법:", False),
            ("\n  할 수 있고 ", False), ("될 수 있는\n  것", True), ("으로 측정", False),
            ("\n\u2022 황금기의 '전체' =\n  ", False), ("모든 사람", True),
        ],
    ],
]


def set_cell_border(cell, color=LIGHT_GRAY):
    """Set thin borders on all sides of a cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = tcPr.find(qn(f"a:{edge}"))
        if ln is not None:
            tcPr.remove(ln)
        ln = tcPr.makeelement(qn(f"a:{edge}"), {})
        ln.set("w", str(Pt(0.5)))
        solidFill = ln.makeelement(qn("a:solidFill"), {})
        srgbClr = solidFill.makeelement(
            qn("a:srgbClr"), {"val": f"{color.rgb:06X}" if hasattr(color, 'rgb') else "CCCCCC"}
        )
        solidFill.append(srgbClr)
        ln.append(solidFill)
        tcPr.append(ln)


def set_cell_fill(cell, color):
    """Set solid fill for a cell."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # Remove existing fill
    existing = tcPr.findall(qn("a:solidFill"))
    for e in existing:
        tcPr.remove(e)
    solidFill = tcPr.makeelement(qn("a:solidFill"), {})
    srgbClr = solidFill.makeelement(
        qn("a:srgbClr"), {"val": f"{color.rgb:06X}" if hasattr(color, 'rgb') else "FFFFFF"}
    )
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def write_cell(cell, fragments, font_size=Pt(8), font_color=DARK_GRAY, line_spacing=1.2):
    """Write rich text fragments to a cell.

    fragments: list of (text, is_bold) tuples.
    """
    cell.text = ""  # clear default
    tf = cell.text_frame
    tf.word_wrap = True

    # Set margins
    tf.margin_left = Mm(2)
    tf.margin_right = Mm(2)
    tf.margin_top = Mm(1.5)
    tf.margin_bottom = Mm(1.5)

    # Vertical alignment
    cell.vertical_anchor = MSO_ANCHOR.TOP

    # Build paragraph(s) from fragments
    # Split on newlines to create separate paragraphs
    paragraphs_data = []
    current_para = []
    for text, bold in fragments:
        parts = text.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                paragraphs_data.append(current_para)
                current_para = []
            if part:
                current_para.append((part, bold))
    if current_para:
        paragraphs_data.append(current_para)

    for pi, para_fragments in enumerate(paragraphs_data):
        if pi == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.alignment = PP_ALIGN.LEFT
        # Set line spacing
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = pPr.makeelement(qn("a:lnSpc"), {})
        spcPct = lnSpc.makeelement(
            qn("a:spcPct"), {"val": str(int(line_spacing * 100000))}
        )
        lnSpc.append(spcPct)
        # Remove existing line spacing
        existing_lnSpc = pPr.findall(qn("a:lnSpc"))
        for e in existing_lnSpc:
            pPr.remove(e)
        pPr.append(lnSpc)

        # Set space after = 0
        spcAft = pPr.makeelement(qn("a:spcAft"), {})
        spcPts = spcAft.makeelement(qn("a:spcPts"), {"val": "0"})
        spcAft.append(spcPts)
        existing_spcAft = pPr.findall(qn("a:spcAft"))
        for e in existing_spcAft:
            pPr.remove(e)
        pPr.append(spcAft)

        for text, bold in para_fragments:
            run = p.add_run()
            run.text = text
            run.font.size = font_size
            run.font.color.rgb = font_color
            run.font.bold = bold
            run.font.name = "Malgun Gothic"


def main():
    prs = Presentation()

    # Set slide size to A3 landscape
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Use blank layout
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # ── Title ──────────────────────────────────────────────────────────────
    from pptx.util import Inches
    title_box = slide.shapes.add_textbox(Mm(15), Mm(8), Mm(390), Mm(16))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "AI 시대 새로운 사회 계약을 위한 기초 연구"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = NAVY
    run.font.name = "Malgun Gothic"

    # ── Subtitle ───────────────────────────────────────────────────────────
    sub_box = slide.shapes.add_textbox(Mm(15), Mm(23), Mm(390), Mm(10))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "태재미래전략연구원 | 2026년 2월"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    run.font.name = "Malgun Gothic"

    # ── Table ──────────────────────────────────────────────────────────────
    n_rows = len(ROWS) + 1  # header + data
    n_cols = len(HEADERS)

    table_left = Mm(10)
    table_top = Mm(36)
    table_width = Mm(400)
    table_height = Mm(248)

    shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top, table_width, table_height)
    table = shape.table

    # Disable default banding
    tbl = table._tbl
    tblPr = tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("bandCol", "0")
    tblPr.set("firstRow", "0")
    tblPr.set("firstCol", "0")
    tblPr.set("lastRow", "0")
    tblPr.set("lastCol", "0")

    # Column widths (proportional)
    col_widths_mm = [38, 62, 58, 62, 62, 62, 56]
    for i, w in enumerate(col_widths_mm):
        table.columns[i].width = Mm(w)

    # Row heights
    header_height = Mm(14)
    # Data rows - distribute remaining space
    data_row_heights = [Mm(42), Mm(56), Mm(56), Mm(80)]

    table.rows[0].height = header_height
    for i, h in enumerate(data_row_heights):
        table.rows[i + 1].height = h

    # ── Header row ─────────────────────────────────────────────────────────
    for ci, header_text in enumerate(HEADERS):
        cell = table.cell(0, ci)
        set_cell_fill(cell, HEADER_BG)
        set_cell_border(cell, RGBColor(0x15, 0x2D, 0x4D))
        write_cell(
            cell,
            [(header_text, True)],
            font_size=Pt(11),
            font_color=WHITE,
            line_spacing=1.1,
        )
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # ── Data rows ──────────────────────────────────────────────────────────
    for ri, row_data in enumerate(ROWS):
        for ci, cell_fragments in enumerate(row_data):
            cell = table.cell(ri + 1, ci)

            # Alternate row backgrounds
            if ri % 2 == 0:
                set_cell_fill(cell, BG_WHITE)
            else:
                set_cell_fill(cell, LIGHT_BG)

            # First column (단계) gets light navy tint
            if ci == 0:
                set_cell_fill(cell, RGBColor(0xE8, 0xEE, 0xF4))

            set_cell_border(cell, LIGHT_GRAY)

            # Dash cells (placeholder) centered and gray
            if len(cell_fragments) == 1 and cell_fragments[0][0] in ("-", ):
                write_cell(
                    cell,
                    [("\u2014", False)],
                    font_size=Pt(9),
                    font_color=RGBColor(0xAA, 0xAA, 0xAA),
                )
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                # Center the dash
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            # Parenthetical cells
            elif len(cell_fragments) == 1 and cell_fragments[0][0].startswith("("):
                write_cell(
                    cell,
                    cell_fragments,
                    font_size=Pt(8),
                    font_color=RGBColor(0x99, 0x99, 0x99),
                )
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                write_cell(cell, cell_fragments, font_size=Pt(8), font_color=DARK_GRAY)

    # ── Footer ─────────────────────────────────────────────────────────────
    footer_box = slide.shapes.add_textbox(Mm(15), Mm(287), Mm(390), Mm(8))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "\u00a9 태재미래전략연구원 2026"
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    run.font.name = "Malgun Gothic"

    # ── Save ───────────────────────────────────────────────────────────────
    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
